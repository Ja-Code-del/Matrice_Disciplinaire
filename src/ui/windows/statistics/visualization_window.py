# src/ui/windows/statistics/visualization_window.py
import os
import openpyxl
from openpyxl.drawing.image import Image
from PyQt6.QtCore import Qt, pyqtSignal
from PyQt6.QtWidgets import (QFileDialog, QDialog, QMessageBox, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QFrame,
                             QLabel,
                             QTableWidget, QPushButton, QTableWidgetItem, QHeaderView, QSizePolicy)
from PyQt6.QtGui import QColor, QIcon
#pour les graphiques et les tableaux
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
import seaborn as sns

#pour excel
import pandas as pd
#pour les pdf
from reportlab.lib import colors, utils
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
#pour powerpoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

from src.data.gendarmerie.structure import SUBDIVISIONS, SERVICE_RANGES, ANALYSIS_THEMES
from src.ui.windows.statistics.chart_selection_dialog import ChartSelectionDialog
from datetime import datetime

from openpyxl.utils import get_column_letter
from openpyxl.styles import Font, PatternFill


class VisualizationWindow(QMainWindow):
    """Fenêtre de visualisation des données statistiques."""

    closed = pyqtSignal()

    def __init__(self, db_manager, config, parent=None):
        super().__init__(parent)
        self.df = None
        self.db_manager = db_manager
        self.config = config
        self.pivot_df = None
        self.current_chart_config = None
        self.setWindowTitle("Visualisation des statistiques")
        self.setMinimumSize(1000, 800)

        # Configuration du style Seaborn
        sns.set_style("whitegrid")
        sns.set_palette("cubehelix", n_colors=8)

        self.setup_ui()
        #self.load_data()

    def show(self):
        """Surcharge de la méthode show pour charger les données avant l'affichage."""
        try:
            # Charger les données avant d'afficher la fenêtre
            self.load_data()
            if hasattr(self, 'df') and self.df is not None:
                super().show()
            else:
                QMessageBox.critical(
                    self,
                    "Erreur",
                    "Impossible de charger les données"
                )
        except Exception as e:
            print(f"Erreur lors de l'affichage: {str(e)}")
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors du chargement des données: {str(e)}"
            )

    def setup_ui(self):
        """Configure l'interface utilisateur."""
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)

        # En-tête avec informations
        header = QFrame()
        header.setFrameStyle(QFrame.Shape.StyledPanel)
        header_layout = QVBoxLayout(header)

        self.title_label = QLabel("Analyse statistique")
        self.title_label.setStyleSheet("font-size: 16px; font-weight: bold;")
        header_layout.addWidget(self.title_label)

        self.info_label = QLabel()
        header_layout.addWidget(self.info_label)
        main_layout.addWidget(header)

        # Tableau de données
        self.table = QTableWidget()
        self.table.setAlternatingRowColors(True)
        self.table.setStyleSheet("""
            QTableWidget {
                background-color: white;
                gridline-color: black;
                border: 1px solid black;
            }
            QTableWidget::item {
                padding: 15px;
            }
        """)
        main_layout.addWidget(self.table)

        # Zone du graphique
        self.figure = Figure(figsize=(8, 6))
        self.canvas = FigureCanvas(self.figure)
        main_layout.addWidget(self.canvas)

        # Boutons d'export
        export_layout = QHBoxLayout()

        self.btn_excel = QPushButton("Exporter Excel")
        self.btn_excel.setStyleSheet("""
                    QPushButton {
                                background-color: #217346;
                                color: white;
                                padding: 8px 15px;
                                border-radius: 15px;
                                font-weight: bold;
                            }
                            QPushButton:hover {
                                background-color: #154734;
                            }
                            """)
        self.btn_excel.setIcon(QIcon("../resources/icons/excel_icon.svg"))
        self.btn_excel.clicked.connect(self.export_excel)

        self.btn_pdf = QPushButton("Exporter PDF")
        self.btn_pdf.setStyleSheet("""
                            QPushButton {
                                        background-color: #6C63FF;
                                        color: white;
                                        padding: 8px 15px;
                                        border-radius: 15px;
                                        font-weight: bold;
                                    }
                                    QPushButton:hover {
                                        background-color: #4B0082;
                                    }
                                    """)
        self.btn_pdf.setIcon(QIcon("../resources/icons/pdf.svg"))
        self.btn_pdf.clicked.connect(self.export_pdf)

        self.btn_pptx = QPushButton("Exporter PowerPoint")
        self.btn_pptx.setStyleSheet("""
                            QPushButton {
                                        background-color: #D24726;
                                        color: white;
                                        padding: 8px 15px;
                                        border-radius: 15px;
                                        font-weight: bold;
                                    }
                                    QPushButton:hover {
                                        background-color: #B7472A;
                                    }
                                    """)

        self.btn_pptx.setIcon(QIcon("../resources/icons/pptx.svg"))
        self.btn_pptx.clicked.connect(self.export_pptx)

        export_layout.addWidget(self.btn_excel)
        export_layout.addWidget(self.btn_pdf)
        export_layout.addWidget(self.btn_pptx)
        export_layout.addStretch()

        main_layout.addLayout(export_layout)

    def get_gendarme_info(self, matricule, numero_dossier):
        """Récupère les informations d'un gendarme pour un dossier spécifique."""
        try:
            gendarme_query = """
            SELECT 
                grade,
                subdiv,
                annee_service,
                situation_matrimoniale
            FROM main_tab 
            WHERE matricule = ? AND numero_dossier = ?
            """
            params = [matricule, numero_dossier]

            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute(gendarme_query, params)
                result = cursor.fetchone()
                return result if result else ("", "", "", "")

        except Exception as e:
            print(f"Erreur dans get_gendarme_info: {str(e)}")
            return ("", "", "", "")

    def load_data(self):
        """Charge et affiche les données selon la configuration."""
        try:
            with self.db_manager.get_connection() as conn:
                subject_selection = self.config.get('subject_selection')
                if not subject_selection:
                    raise Exception("Aucun sujet d'analyse sélectionné")

                # Requête de base pour obtenir les sanctions
                sanctions_query = """
                SELECT 
                    id,
                    matricule,
                    numero_dossier,
                    date_enr,
                    date_faits,
                    faute_commise,
                    categorie,
                    statut,
                    annee_punition,
                    annee_faits
                FROM main_tab
                """

                params = []
                if not subject_selection['value'].startswith('Tous'):
                    field = subject_selection['field']
                    value = subject_selection['value']
                    sanctions_query += f" WHERE {field} = ?"
                    params.append(value)

                # Exécuter la requête pour obtenir les sanctions
                sanctions_df = pd.read_sql_query(sanctions_query, conn, params=params)

                # Créer une liste pour stocker toutes les lignes de données
                data_rows = []

                # Pour chaque sanction, récupérer les infos du gendarme
                for _, row in sanctions_df.iterrows():
                    gendarme_info = self.get_gendarme_info(row['matricule'], row['numero_dossier'])

                    data_row = {
                        'id': row['id'],
                        'matricule': row['matricule'],
                        'numero_dossier': row['numero_dossier'],
                        'date_enr': row['date_enr'],
                        'date_faits': row['date_faits'],
                        'faute_commise': row['faute_commise'],
                        'categorie': row['categorie'],
                        'statut': row['statut'],
                        'annee_punition': row['annee_punition'],
                        'annee_faits': row['annee_faits'],
                        'grade': gendarme_info[0],
                        'subdiv': gendarme_info[1],
                        'annee_service': gendarme_info[2],
                        'situation_matrimoniale': gendarme_info[3]
                    }
                    data_rows.append(data_row)

                # Créer le DataFrame final
                df = pd.DataFrame(data_rows)

                if df.empty:
                    raise Exception("Aucune donnée disponible")

                # Suite du code pour le traitement des années de service et la création des graphiques
                # [Le reste du code reste identique]

                # Traitement des années de service
                if any(config["field"] == "annee_service" for config in
                       [self.config["x_axis"], self.config["y_axis"]]):
                    service_categories = ['0-5 ANS', '6-10 ANS', '11-15 ANS', '16-20 ANS',
                                          '21-25 ANS', '26-30 ANS', '31-35 ANS', '36-40 ANS', 'Non spécifié']

                    df['annee_service'] = pd.to_numeric(df['annee_service'], errors='coerce')

                    df['service_range'] = pd.cut(
                        df['annee_service'],
                        bins=[0, 5, 10, 15, 20, 25, 30, 35, 40],
                        labels=service_categories[:-1],
                        include_lowest=True
                    )

                    df['service_range'] = df['service_range'].cat.add_categories(['Non spécifié'])
                    df.loc[df['service_range'].isna(), 'service_range'] = 'Non spécifié'

                # Préparation des données pour les axes
                x_config = self.config["x_axis"]
                y_config = self.config["y_axis"]

                x_values = df['service_range'] if x_config["field"] == "annee_service" else df[x_config["field"]]
                y_values = df['service_range'] if y_config["field"] == "annee_service" else df[y_config["field"]]

                pivot_df = pd.crosstab(
                    index=y_values,
                    columns=x_values,
                    dropna=True,
                    margins=True,
                    margins_name='TOTAL'
                )

                graph_df = pd.DataFrame({
                    'x_value': x_values,
                    'y_value': y_values
                })
                graph_df = graph_df.groupby(['x_value', 'y_value'], observed=True).size().reset_index(name='count')

                self.df = graph_df
                self.pivot_df = pivot_df
                self.update_table(pivot_df)

        except Exception as e:
            print(f"Error in load_data: {str(e)}")
            import traceback
            traceback.print_exc()
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors du chargement des données: {str(e)}"
            )

    def cleanup(self):
        plt.close('all')
        if hasattr(self, 'canvas'):
            self.canvas.close()


    def apply_filters(self, df):
        """Applique les filtres de configuration au DataFrame."""
        try:
            filtered_df = df.copy()

            # Application des filtres pour chaque axe
            for config in [self.config["x_axis"], self.config["y_axis"]]:
                value = config.get("value", "")
                if value and not value.startswith("Tous"):
                    if config["field"] == "annee_service":
                        service_range = value
                        start, end = map(int, service_range.split("-")[0:2])
                        filtered_df = filtered_df[
                            (filtered_df['annee_service'] >= start) &
                            (filtered_df['annee_service'] <= end)
                            ]
                    else:
                        field = config["field"]
                        filtered_df = filtered_df[filtered_df[field] == value]

            return filtered_df

        except Exception as e:
            print(f"Error in apply_filters: {str(e)}")
            return df  # Retourner le DataFrame original en cas d'erreur

    def update_table(self, pivot_df):
        """Met à jour le tableau avec les données sous forme de tableau croisé."""
        try:
            print("Tableau pivot:", pivot_df)  # Debug

            # Vérifier si on a des données
            if pivot_df is None:
                return

            # Configuration du tableau
            self.table.clear()
            self.table.setRowCount(len(pivot_df.index))
            self.table.setColumnCount(len(pivot_df.columns))

            # En-têtes
            self.table.setHorizontalHeaderLabels(pivot_df.columns.astype(str))
            self.table.setVerticalHeaderLabels(pivot_df.index.astype(str))

            # Remplissage des données
            for i in range(len(pivot_df.index)):
                for j in range(len(pivot_df.columns)):
                    try:
                        value = pivot_df.iloc[i, j]
                        # S'assurer que la valeur est numérique
                        if pd.isna(value):
                            str_value = "0"
                        else:
                            str_value = str(int(value))
                    except:
                        str_value = "0"

                    item = QTableWidgetItem(str_value)
                    item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)

                    # Style pour les totaux
                    if (pivot_df.index[i] == 'TOTAL' or
                            pivot_df.columns[j] == 'TOTAL'):
                        item.setBackground(QColor(128, 128, 128))
                        item.setForeground(QColor(255, 255, 255))
                        font = item.font()
                        font.setBold(True)
                        item.setFont(font)
                    else:
                        item.setBackground(QColor(220, 220, 220))

                    self.table.setItem(i, j, item)

            # Style des en-têtes
            header_style = """
                QHeaderView::section {
                    background-color: rgb(0, 85, 127);
                    color: white;
                    padding: 6px;
                    border: 1px solid #005580;
                    font-weight: bold;
                }
            """
            self.table.horizontalHeader().setStyleSheet(header_style)
            self.table.verticalHeader().setStyleSheet(header_style)

            # Ajustement des dimensions
            self.table.horizontalHeader().setSectionResizeMode(QHeaderView.ResizeMode.Stretch)
            self.table.verticalHeader().setSectionResizeMode(QHeaderView.ResizeMode.Stretch)

        except Exception as e:
            print(f"Erreur dans update_table: {str(e)}")
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de la mise à jour du tableau: {str(e)}"
            )

    def update_graph(self, df, chart_config):
        """Met à jour le graphique selon la configuration choisie."""
        try:
            # Vérifier les données et la configuration
            if self.df is None or self.pivot_df is None:
                raise Exception("Aucune donnée disponible")

            if not isinstance(chart_config, dict):
                raise Exception("Configuration du graphique invalide")

            # Nettoyage du graphique précédent
            self.figure.clear()
            ax = self.figure.add_subplot(111)

            # Extraire les informations de configuration
            chart_type = chart_config.get('type')
            selected_axis = chart_config.get('axis', 'x')

            # Configuration du titre
            subject = self.config['subject_selection']
            subject_title = f"{subject['theme']}: {subject['value']}"

            # Choix du graphique selon le type
            if chart_type == 'bar_simple':
                if chart_config.get('orientation') == 'horizontal':
                    self._create_simple_bar_horizontal_sns(ax, selected_axis, subject_title)
                else:
                    self._create_simple_bar_sns(ax, selected_axis, subject_title)
            elif chart_type == 'bar_grouped':
                if chart_config.get('orientation') == 'horizontal':
                    self._create_bar_grouped_horizontal_sns(ax, subject_title)
                else:
                    self._create_bar_grouped_sns(ax, subject_title)
            elif chart_type == 'pie':
                self._create_pie_chart_sns(ax, selected_axis, subject_title)
            elif chart_type == 'donut':
                self._create_donut_chart_sns(ax, selected_axis, subject_title)
            elif chart_type == 'bar_stacked':
                self._create_stacked_bar_sns(ax, subject_title)
            elif chart_type == 'heatmap':
                self._create_heatmap_sns(ax, subject_title)
            else:
                self._create_stacked_bar_sns(ax, subject_title)  # Par défaut

            # Ajustement final du graphique
            self.figure.tight_layout()
            self.canvas.draw()

        except Exception as e:
            print(f"Erreur dans update_graph: {str(e)}")
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de la mise à jour du graphique: {str(e)}"
            )

    def update_info(self, df):
        """Met à jour les informations d'en-tête."""
        try:
            # Vérifier si on a des données
            if df is None:
                return

            total = len(df)
            moyenne = df['annee_punition'].mean() if 'annee_punition' in df.columns else 0
            maximum = df['annee_punition'].max() if 'annee_punition' in df.columns else 0
            minimum = df['annee_punition'].min() if 'annee_punition' in df.columns else 0

            self.info_label.setText(
                f"Total: {total} | "
                f"Moyenne: {moyenne:.2f} | "
                f"Maximum: {maximum} | "
                f"Minimum: {minimum}"
            )
        except Exception as e:
            print(f"Erreur dans update_info: {str(e)}")

    #Methodes de créations des graphiques
    def _create_simple_bar_sns(self, ax, selected_axis, subject_title):
        """Crée un histogramme simple selon l'axe choisi."""
        try:
            # Sélectionner les données selon l'axe
            value_col = 'x_value' if selected_axis == 'x' else 'y_value'
            config = self.config[f'{selected_axis}_axis']

            # Agréger les données
            data = self.df.groupby(value_col)['count'].sum().reset_index()

            # Supprimer les lignes avec compte = 0 ou valeur = "Non spécifié"
            data = data[
                (data['count'] > 0) &
                (data[value_col] != "Non spécifié")
                ]

            # Trier les données par valeur décroissante
            data = data.sort_values('count', ascending=False)

            # Création du graphique
            sns.barplot(
                data=data,
                x=value_col,
                y='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=len(data))
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\nDistribution par {config['theme']}")
            ax.set_xlabel(config['theme'])
            ax.set_ylabel('Nombre de sanctions')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout des valeurs sur les barres
            for i, v in enumerate(data['count']):
                ax.text(
                    i, v + 0.5,
                    str(int(v)),
                    ha='center',
                    va='bottom',
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_simple_bar_sns: {str(e)}")
            raise

    def _create_pie_chart_sns(self, ax, selected_axis, subject_title):
        """Crée un graphique en camembert selon l'axe choisi."""
        try:
            # Sélectionner les données selon l'axe
            value_col = 'x_value' if selected_axis == 'x' else 'y_value'
            config = self.config[f'{selected_axis}_axis']

            # Agréger les données
            data = self.df.groupby(value_col)['count'].sum()

            # Calculer les pourcentages
            total = data.sum()

            # Création du graphique
            wedges, texts, autotexts = ax.pie(
                data,
                labels=[f"{x}\n({v:,} - {(v / total) * 100:.1f}%)" for x, v in data.items()],
                autopct='',
                colors=sns.color_palette("cubehelix", n_colors=len(data))
            )

            # Titre
            ax.set_title(f"{subject_title}\nRépartition par {config['theme']}")

            # Légende
            ax.legend(
                wedges,
                [f"{x}" for x in data.index],
                title=config['theme'],
                loc="center left",
                bbox_to_anchor=(1, 0, 0.5, 1)
            )

        except Exception as e:
            print(f"Erreur dans _create_pie_chart_sns: {str(e)}")
            raise

    def _create_donut_chart_sns(self, ax, selected_axis, subject_title):
        """Crée un graphique en anneau selon l'axe choisi."""
        try:
            # Sélectionner les données selon l'axe
            value_col = 'x_value' if selected_axis == 'x' else 'y_value'
            config = self.config[f'{selected_axis}_axis']

            # Agréger les données
            data = self.df.groupby(value_col)['count'].sum()

            # Calculer les pourcentages et le total
            total = data.sum()

            # Création du graphique
            wedges, texts, autotexts = ax.pie(
                data,
                labels=[f"{x}\n({v:,} - {(v / total) * 100:.1f}%)" for x, v in data.items()],
                autopct='',
                colors=sns.color_palette("cubehelix", n_colors=len(data)),
                wedgeprops=dict(width=0.5)  # Crée l'effet donut
            )

            # Ajouter le total au centre
            centre_circle = plt.Circle((0, 0), 0.70, fc='white')
            ax.add_artist(centre_circle)
            ax.text(0, 0, f'Total\n{total:,}', ha='center', va='center', fontsize=12, fontweight='bold')

            # Titre
            ax.set_title(f"{subject_title}\nRépartition par {config['theme']}")

            # Légende
            ax.legend(
                wedges,
                [f"{x}" for x in data.index],
                title=config['theme'],
                loc="center left",
                bbox_to_anchor=(1, 0, 0.5, 1)
            )

        except Exception as e:
            print(f"Erreur dans _create_donut_chart_sns: {str(e)}")
            raise

    def _create_stacked_bar_sns(self, ax, subject_title):
        """Crée un graphique à barres empilées."""
        try:
            # Utilisation directe du pivot_df pour les barres empilées
            # Retirer la ligne 'TOTAL' pour le graphique
            plot_df = self.pivot_df.drop('TOTAL', axis=0).drop('TOTAL', axis=1)

            # Création du graphique
            plot_df.plot(
                kind='bar',
                stacked=True,
                ax=ax,
                colormap='cubehelix'
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\n"
                         f"Répartition {self.config['y_axis']['theme']} "
                         f"par {self.config['x_axis']['theme']}")
            ax.set_xlabel(self.config['x_axis']['theme'])
            ax.set_ylabel('Nombre de sanctions')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout des totaux au-dessus de chaque barre
            totals = plot_df.sum(axis=1)
            for i, total in enumerate(totals):
                ax.text(
                    i, total, f'{int(total)}',
                    ha='center',
                    va='bottom',
                    fontweight='bold'
                )

            # Légende
            ax.legend(
                title=self.config['y_axis']['theme'],
                bbox_to_anchor=(1.05, 1),
                loc='upper left'
            )

        except Exception as e:
            print(f"Erreur dans _create_stacked_bar_sns: {str(e)}")
            raise

    def _create_heatmap_sns(self, ax, subject_title):
        """Crée une heatmap des données."""
        try:
            # Utiliser le pivot_df sans les totaux
            plot_df = self.pivot_df.drop('TOTAL', axis=0).drop('TOTAL', axis=1)

            # Création de la heatmap
            sns.heatmap(
                plot_df,
                annot=True,
                fmt='d',
                cmap='YlOrRd',
                ax=ax,
                cbar_kws={'label': 'Nombre de sanctions'}
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\n"
                         f"Distribution {self.config['y_axis']['theme']} "
                         f"vs {self.config['x_axis']['theme']}")

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')
            plt.yticks(rotation=0)

            # Ajustement de la taille des étiquettes si nécessaire
            ax.set_xticklabels(ax.get_xticklabels(), fontsize=8)
            ax.set_yticklabels(ax.get_yticklabels(), fontsize=8)

        except Exception as e:
            print(f"Erreur dans _create_heatmap_sns: {str(e)}")
            raise

    def _create_bar_grouped_sns(self, ax, subject_title):
        """Crée un graphique à barres groupées."""
        try:
            # Réorganiser les données pour le graphique groupé
            df_melted = self.df.copy()

            # Définition des couleurs selon le type de sous-thème
            if self.config['y_axis']['theme'] == "Situation matrimoniale":
                colors = {'CELIBATAIRE': '#0066cc', 'MARIE': '#ff9933'}
            else:
                # Générer une palette de couleurs selon le nombre de valeurs uniques
                unique_values = df_melted['y_value'].unique()
                n_colors = len(unique_values)
                colors = {}
                color_palette = sns.color_palette("husl", n_colors)  # Utiliser husl pour une meilleure distinction
                for i, value in enumerate(unique_values):
                    colors[value] = color_palette[i]

            print("Valeurs uniques:", df_melted['y_value'].unique())  # Debug
            print("Palette de couleurs:", colors)  # Debug

            # Création du graphique
            sns.barplot(
                data=df_melted,
                x='x_value',
                y='count',
                hue='y_value',
                ax=ax,
                palette=colors
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\n"
                         f"Comparaison {self.config['y_axis']['theme']} "
                         f"par {self.config['x_axis']['theme']}")
            ax.set_xlabel(self.config['x_axis']['theme'])
            ax.set_ylabel('Nombre de sanctions')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout des valeurs sur les barres
            for container in ax.containers:
                ax.bar_label(container, padding=3)

            # Ajuster la légende selon le nombre d'éléments
            if len(unique_values) > 6:  # Si beaucoup de valeurs
                ax.legend(
                    title=self.config['y_axis']['theme'],
                    bbox_to_anchor=(1.05, 1),
                    loc='upper left',
                    ncol=max(1, len(unique_values) // 8)  # Diviser en colonnes si nécessaire
                )
            else:
                ax.legend(
                    title=self.config['y_axis']['theme'],
                    bbox_to_anchor=(1.05, 1),
                    loc='upper left'
                )

            # Ajustement de la figure pour éviter le chevauchement
            plt.tight_layout()

        except Exception as e:
            print(f"Erreur dans _create_bar_grouped_sns: {str(e)}")
            raise

    def _create_line_chart_sns(self, ax, subject_title):
        """Crée un graphique linéaire."""
        try:
            # Utiliser le pivot_df sans les totaux
            plot_df = self.pivot_df.drop('TOTAL', axis=0).drop('TOTAL', axis=1)

            # Tracer une ligne pour chaque colonne
            for column in plot_df.columns:
                ax.plot(
                    plot_df.index,
                    plot_df[column],
                    marker='o',
                    label=column,
                    linewidth=2,
                    markersize=8
                )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\n"
                         f"Évolution {self.config['y_axis']['theme']} "
                         f"par {self.config['x_axis']['theme']}")
            ax.set_xlabel(self.config['x_axis']['theme'])
            ax.set_ylabel('Nombre de sanctions')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout des valeurs sur les points
            for line in ax.lines:
                for x, y in zip(range(len(line.get_xdata())), line.get_ydata()):
                    ax.text(
                        x, y,
                        f'{int(y)}',
                        ha='center',
                        va='bottom'
                    )

            # Grille
            ax.grid(True, linestyle='--', alpha=0.7)

            # Légende
            ax.legend(
                title=self.config['y_axis']['theme'],
                bbox_to_anchor=(1.05, 1),
                loc='upper left'
            )

        except Exception as e:
            print(f"Erreur dans _create_line_chart_sns: {str(e)}")
            raise

    def _create_box_plot_sns(self, ax, df, axis='x'):
        """Crée une boîte à moustaches avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique
            sns.boxplot(
                data=df,
                x=group_col,
                y='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=8)
            )

            # Style et étiquettes
            ax.set_title(f"Distribution statistique par {title_theme}")
            ax.set_xlabel('')
            ax.set_ylabel('Nombre')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout de la grille
            ax.grid(True, linestyle='--', alpha=0.7)

            # Ajout des statistiques sur chaque boîte
            for i, artist in enumerate(ax.artists):
                stats = df[df[group_col] == df[group_col].unique()[i]]['count'].describe()
                ax.text(
                    i, stats['max'],
                    f'Max: {int(stats["max"]):,}\n'
                    f'Méd: {int(stats["50%"]):,}\n'
                    f'Min: {int(stats["min"]):,}',
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_box_plot_sns: {str(e)}")
            raise

    def _create_violin_plot_sns(self, ax, df, axis='x'):
        """Crée un violin plot avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique
            sns.violinplot(
                data=df,
                x=group_col,
                y='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=8),
                inner='box'  # Affiche un boxplot à l'intérieur
            )

            # Style et étiquettes
            ax.set_title(f"Distribution détaillée par {title_theme}")
            ax.set_xlabel('')
            ax.set_ylabel('Nombre')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout de la grille
            ax.grid(True, linestyle='--', alpha=0.7)

            # Ajout des statistiques
            for i, label in enumerate(df[group_col].unique()):
                stats = df[df[group_col] == label]['count'].describe()
                ax.text(
                    i, stats['max'],
                    f'Méd: {int(stats["50%"]):,}',
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_violin_plot_sns: {str(e)}")
            raise

    def _create_kde_plot_sns(self, ax, df, axis='x'):
        """Crée un KDE plot (estimation par noyau de la densité) avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique pour chaque catégorie
            categories = df[group_col].unique()
            palette = sns.color_palette("cubehelix", n_colors=len(categories))

            for idx, category in enumerate(categories):
                subset = df[df[group_col] == category]
                sns.kdeplot(
                    data=subset,
                    x='count',
                    label=category,
                    ax=ax,
                    color=palette[idx],
                    fill=True,
                    alpha=0.5
                )

            # Style et étiquettes
            ax.set_title(f"Densité de distribution par {title_theme}")
            ax.set_xlabel('Nombre')
            ax.set_ylabel('Densité')

            # Légende
            ax.legend(
                title=title_theme,
                bbox_to_anchor=(1.05, 1),
                loc='upper left'
            )

            # Grille
            ax.grid(True, linestyle='--', alpha=0.7)

        except Exception as e:
            print(f"Erreur dans _create_kde_plot_sns: {str(e)}")
            raise

    def _create_swarm_plot_sns(self, ax, df, axis='x'):
        """Crée un swarm plot (nuage de points) avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique
            sns.swarmplot(
                data=df,
                x=group_col,
                y='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=8),
                size=8,
                alpha=0.7
            )

            # Style et étiquettes
            ax.set_title(f"Distribution détaillée par {title_theme}")
            ax.set_xlabel('')
            ax.set_ylabel('Nombre')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout de la grille
            ax.grid(True, linestyle='--', alpha=0.7)

            # Ajout des statistiques pour chaque groupe
            for i, label in enumerate(df[group_col].unique()):
                stats = df[df[group_col] == label]['count'].describe()
                ax.text(
                    i, stats['max'],
                    f'Moy: {int(stats["mean"]):,}\n'
                    f'N: {int(stats["count"]):,}',
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_swarm_plot_sns: {str(e)}")
            raise

    def _create_count_plot_sns(self, ax, df, axis='x'):
        """Crée un count plot (histogramme de fréquences) avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique
            sns.countplot(
                data=df,
                x=group_col,
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=8),
                order=df[group_col].value_counts().index
            )

            # Style et étiquettes
            ax.set_title(f"Fréquence des occurrences par {title_theme}")
            ax.set_xlabel('')
            ax.set_ylabel('Nombre d\'occurrences')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout de la grille
            ax.grid(True, axis='y', linestyle='--', alpha=0.7)

            # Ajout des valeurs sur les barres
            for i, v in enumerate(df[group_col].value_counts()):
                ax.text(
                    i, v,
                    f'{int(v):,}',
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_count_plot_sns: {str(e)}")
            raise

    def _create_strip_plot_sns(self, ax, df, axis='x'):
        """Crée un strip plot (nuage de points en ligne) avec Seaborn."""
        try:
            # Sélectionner la colonne appropriée selon l'axe choisi
            group_col = 'x_value' if axis == 'x' else 'y_value'
            title_theme = self.config[f'{axis}_axis']['theme']

            # Création du graphique
            sns.stripplot(
                data=df,
                x=group_col,
                y='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=8),
                size=8,
                alpha=0.6,
                jitter=True  # Ajoute un peu de bruit aléatoire pour éviter la superposition
            )

            # Style et étiquettes
            ax.set_title(f"Distribution points par {title_theme}")
            ax.set_xlabel('')
            ax.set_ylabel('Nombre')

            # Rotation des étiquettes
            plt.xticks(rotation=45, ha='right')

            # Ajout de la grille
            ax.grid(True, linestyle='--', alpha=0.7)

            # Ajout des statistiques pour chaque groupe
            for i, label in enumerate(df[group_col].unique()):
                stats = df[df[group_col] == label]['count'].describe()
                ax.text(
                    i, stats['max'],
                    f'Moy: {int(stats["mean"]):,}\n'
                    f'Méd: {int(stats["50%"]):,}',
                    ha='center',
                    va='bottom',
                    fontsize=8,
                    fontweight='bold'
                )

            # Ajout d'une ligne de moyenne pour chaque groupe
            for i, label in enumerate(df[group_col].unique()):
                mean_val = df[df[group_col] == label]['count'].mean()
                ax.hlines(
                    y=mean_val,
                    xmin=i - 0.2,
                    xmax=i + 0.2,
                    color='red',
                    linestyles='dashed',
                    alpha=0.5,
                    label='Moyenne' if i == 0 else ""
                )

            # Légende pour la ligne de moyenne
            if len(df[group_col].unique()) > 0:
                ax.legend(['Moyenne'], loc='upper right')

        except Exception as e:
            print(f"Erreur dans _create_strip_plot_sns: {str(e)}")
            raise

    def _create_simple_bar_horizontal_sns(self, ax, selected_axis, subject_title):
        """Crée un histogramme simple horizontal selon l'axe choisi."""
        try:
            # Sélectionner les données selon l'axe
            value_col = 'x_value' if selected_axis == 'x' else 'y_value'
            config = self.config[f'{selected_axis}_axis']

            # Agréger les données
            data = self.df.groupby(value_col)['count'].sum().reset_index()

            # Trier les données par valeur décroissante
            data = data.sort_values('count', ascending=True)  # Ascending True pour l'affichage horizontal

            # Création du graphique
            sns.barplot(
                data=data,
                y=value_col,  # Inverser x et y pour l'orientation horizontale
                x='count',
                ax=ax,
                palette=sns.color_palette("cubehelix", n_colors=len(data)),
                orient='h'  # Spécifier l'orientation horizontale
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\nDistribution par {config['theme']}")
            ax.set_ylabel(config['theme'])
            ax.set_xlabel('Nombre de sanctions')

            # Pas besoin de rotation pour les étiquettes verticales
            ax.tick_params(axis='y', labelrotation=0)

            # Ajout des valeurs sur les barres
            for i, v in enumerate(data['count']):
                ax.text(
                    v + 0.5,  # Décalage horizontal
                    i,  # Position verticale
                    str(int(v)),
                    ha='left',
                    va='center',
                    fontweight='bold'
                )

        except Exception as e:
            print(f"Erreur dans _create_simple_bar_horizontal_sns: {str(e)}")
            raise

    def _create_bar_grouped_horizontal_sns(self, ax, subject_title):
        """Crée un graphique à barres groupées horizontal."""
        try:
            # Réorganiser les données pour le graphique groupé
            df_melted = self.df.copy()

            # Création du graphique
            sns.barplot(
                data=df_melted,
                y='x_value',  # Inverser x et y pour l'orientation horizontale
                x='count',
                hue='y_value',
                ax=ax,
                palette='cubehelix',
                orient='h'  # Spécifier l'orientation horizontale
            )

            # Style et étiquettes
            ax.set_title(f"{subject_title}\n"
                         f"Comparaison {self.config['y_axis']['theme']} "
                         f"par {self.config['x_axis']['theme']}")
            ax.set_ylabel(self.config['x_axis']['theme'])
            ax.set_xlabel('Nombre de sanctions')

            # Pas besoin de rotation pour les étiquettes verticales
            ax.tick_params(axis='y', labelrotation=0)

            # Ajout des valeurs sur les barres
            for container in ax.containers:
                ax.bar_label(
                    container,
                    padding=3,
                    label_type='edge'  # Placer les labels à l'extrémité des barres
                )

            # Légende
            ax.legend(
                title=self.config['y_axis']['theme'],
                bbox_to_anchor=(1.05, 1),
                loc='upper left'
            )

            # Ajustement de la figure pour éviter le chevauchement
            plt.tight_layout()

        except Exception as e:
            print(f"Erreur dans _create_bar_grouped_horizontal_sns: {str(e)}")
            raise

    def closeEvent(self, event):
        """Gère la fermeture propre de la fenêtre."""
        if hasattr(self, 'figure'):
            self.figure.clear()
        self.closed.emit()
        event.accept()

    # Méthodes d'export
    def export_excel(self):
        """Exporte les données vers Excel avec headers de ligne et de colonne."""
        try:
            default_name = f"statistiques_excel_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer le fichier Excel",
                default_name,
                "Excel Files (*.xlsx)"
            )

            if not file_path:
                return

            with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
                # Si nous avons un pivot_df, l'utiliser directement
                if hasattr(self, 'pivot_df'):
                    # Export du DataFrame pivot avec ses index
                    self.pivot_df.to_excel(writer, sheet_name='Données')
                    worksheet = writer.sheets['Données']

                    # Mise en forme des colonnes
                    for idx, col in enumerate(self.pivot_df.columns, start=1):
                        col_letter = get_column_letter(idx + 1)  # +1 car la première colonne contient les index
                        max_length = max(
                            self.pivot_df[col].astype(str).apply(len).max(),
                            len(str(col))
                        ) + 2
                        worksheet.column_dimensions[col_letter].width = max_length

                    # Mise en forme de la colonne d'index
                    max_index_length = max(
                        max(len(str(idx)) for idx in self.pivot_df.index),
                        len(str(self.pivot_df.index.name)) if self.pivot_df.index.name else 0
                    ) + 2
                    worksheet.column_dimensions['A'].width = max_index_length

                    # Style pour les totaux (lignes et colonnes)
                    for row in worksheet.iter_rows():
                        for cell in row:
                            if 'TOTAL' in str(cell.value):
                                cell.font = Font(bold=True)
                                cell.fill = PatternFill(
                                    start_color="808080",
                                    end_color="808080",
                                    fill_type="solid"
                                )

                # Export des données du graphique si disponibles
                if hasattr(self, 'df'):
                    self.df.to_excel(writer, sheet_name='Données Graphique', index=False)
                    graph_sheet = writer.sheets['Données Graphique']

                    # Mise en forme de la feuille graphique
                    for idx, col in enumerate(self.df.columns, start=1):
                        max_length = max(
                            self.df[col].astype(str).apply(len).max(),
                            len(str(col))
                        ) + 2
                        graph_sheet.column_dimensions[get_column_letter(idx)].width = max_length

                    # Feuille de configuration
                    info_df = pd.DataFrame({
                        'Configuration': [
                            'Axe X',
                            'Axe Y',
                            'Sélection',
                            'Type de graphique'
                        ],
                        'Valeur': [
                            f"{self.config['x_axis']['field']}",
                            f"{self.config['y_axis']['field']}",
                            f"{self.config['subject_selection']['field']} - {self.config['subject_selection']['value']}",
                            self.config.get('graph_type', 'Non spécifié')
                        ]
                    })
                    info_df.to_excel(writer, sheet_name='Configuration', index=False)

                    # Ajustement des colonnes de la feuille de configuration
                    config_sheet = writer.sheets['Configuration']
                    for column in config_sheet.columns:
                        max_length = 0
                        for cell in column:
                            try:
                                max_length = max(max_length, len(str(cell.value)))
                            except:
                                pass
                        config_sheet.column_dimensions[get_column_letter(column[0].column)].width = max_length + 2

                        # Ajouter une nouvelle feuille pour le graphique
                        if hasattr(self, 'graph_widget'):
                            workbook = writer.book
                            graph_sheet = workbook.create_sheet("Graphique")

                            # Ajouter le graphique
                            img = openpyxl.drawing.image.Image(temp_img_path)
                            img.anchor = 'B2'  # Position du graphique
                            graph_sheet.add_image(img)

                            # Ajuster la hauteur de la ligne et la largeur de la colonne
                            graph_sheet.row_dimensions[2].height = 300
                            graph_sheet.column_dimensions['B'].width = 60

            # Supprimer le fichier temporaire
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            print(f"Erreur détaillée: {str(e)}")
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export Excel:\n{str(e)}"
            )



    def export_pdf(self):
        """Exporte les données vers PDF avec mise en forme."""
        try:
            # Demander à l'utilisateur où sauvegarder le fichier
            default_name = f"statistiques_pdf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer le fichier PDF",
                default_name,
                "PDF Files (*.pdf)"
            )

            if not file_path:  # L'utilisateur a annulé
                return

            # Création du document PDF
            doc = SimpleDocTemplate(
                file_path,
                pagesize=landscape(letter),
                rightMargin=30,
                leftMargin=30,
                topMargin=30,
                bottomMargin=30
            )

            # Liste des éléments à ajouter au PDF
            elements = []
            styles = getSampleStyleSheet()

            # Titre du document
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=30,
                alignment=1  # Centre
            )
            title = Paragraph("Rapport Statistique", title_style)
            elements.append(title)

            # Date et informations
            date_style = ParagraphStyle(
                'DateStyle',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=20
            )
            date_text = Paragraph(
                f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}",
                date_style
            )
            elements.append(date_text)

            # Informations de configuration
            config_style = ParagraphStyle(
                'ConfigStyle',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=20
            )
            config_text = Paragraph(
                f"Analyse: {self.config['x_axis']['theme']} par {self.config['y_axis']['theme']}",
                config_style
            )
            elements.append(config_text)
            elements.append(Spacer(1, 20))

            # Statistiques globales
            stats_text = self.info_label.text()
            stats_para = Paragraph(stats_text, styles['Normal'])
            elements.append(stats_para)
            elements.append(Spacer(1, 20))

            # Création des données du tableau
            table_data = []
            headers = []

            # En-têtes
            for j in range(self.table.columnCount()):
                headers.append(self.table.horizontalHeaderItem(j).text())
            table_data.append(headers)

            # Données
            for i in range(self.table.rowCount()):
                row = []
                for j in range(self.table.columnCount()):
                    item = self.table.item(i, j)
                    row.append(item.text() if item else "")
                table_data.append(row)

            # Style du tableau
            table_style = TableStyle([
                # En-têtes
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#005580')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                # Corps du tableau
                ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                # Ligne totaux
                ('BACKGROUND', (-1, -1), (-1, -1), colors.grey),
                ('TEXTCOLOR', (-1, -1), (-1, -1), colors.whitesmoke),
                ('FONTNAME', (-1, -1), (-1, -1), 'Helvetica-Bold'),
            ])

            # Création du tableau PDF
            pdf_table = Table(table_data, repeatRows=1)
            pdf_table.setStyle(table_style)
            elements.append(pdf_table)

            # Ajout du graphique
            elements.append(Spacer(1, 30))

            # Sauvegarde temporaire du graphique
            graph_path = 'temp_graph.png'
            self.figure.savefig(graph_path, dpi=300, bbox_inches='tight')

            # Ajout de l'image au PDF
            img = Image(graph_path, width=500, height=300)
            elements.append(img)

            # Construction du document
            doc.build(elements)

            # Nettoyage du fichier temporaire
            if os.path.exists(graph_path):
                os.remove(graph_path)

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export PDF:\n{str(e)}"
            )

    def export_pptx(self):
        """Exporte les données vers PowerPoint avec mise en forme."""
        try:
            # Demander à l'utilisateur où sauvegarder le fichier
            default_name = f"statistiques_pptx_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer la présentation PowerPoint",
                default_name,
                "PowerPoint Files (*.pptx)"
            )

            if not file_path:  # L'utilisateur a annulé
                return

            # Création de la présentation
            prs = Presentation()

            # Définir la taille des diapositives (16:9)
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Slide de titre
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Rapport Statistique"
            subtitle.text = f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}"

            # Slide d'information
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = "Informations de l'analyse"

            tf = body_shape.text_frame
            tf.text = f"Analyse: {self.config['x_axis']['theme']} par {self.config['y_axis']['theme']}"

            p = tf.add_paragraph()
            p.text = self.info_label.text()

            # Slide avec le tableau
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            title = shapes.title
            title.text = "Données"

            # Création du tableau
            rows = self.table.rowCount()
            cols = self.table.columnCount()

            left = Inches(0.5)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1)
            height = Inches(0.3 * rows)

            table = shapes.add_table(rows, cols, left, top, width, height).table

            # Style de l'en-tête
            for i in range(cols):
                cell = table.cell(0, i)
                cell.text = self.table.horizontalHeaderItem(i).text()

                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER

                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 85, 128)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # Remplissage des données
            for i in range(1, rows):
                for j in range(cols):
                    cell = table.cell(i, j)
                    item = self.table.item(i - 1, j)
                    cell.text = item.text() if item else ""

                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER

                    # Style pour la ligne des totaux
                    if i == rows - 1:
                        paragraph.font.bold = True
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)

            # Slide avec le graphique
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            title = shapes.title
            title.text = "Visualisation"

            # Sauvegarde temporaire du graphique
            graph_path = 'temp_graph.png'
            self.figure.savefig(graph_path, dpi=300, bbox_inches='tight')

            # Ajout du graphique
            left = Inches(1)
            top = Inches(1.5)
            pic = slide.shapes.add_picture(
                graph_path,
                left,
                top,
                width=prs.slide_width - Inches(2)
            )

            # Sauvegarde de la présentation
            prs.save(file_path)

            # Nettoyage du fichier temporaire
            if os.path.exists(graph_path):
                os.remove(graph_path)

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export PowerPoint:\n{str(e)}"
            )
