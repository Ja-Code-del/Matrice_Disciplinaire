# src/ui/windows/statistics/full_list_window.py

from PyQt6.QtWidgets import (QMainWindow, QWidget, QFileDialog, QVBoxLayout, QHBoxLayout,
                             QPushButton, QTableWidget, QTableWidgetItem,
                             QLabel, QComboBox, QLineEdit, QFrame,
                             QHeaderView, QMessageBox, QGroupBox, QRadioButton, QButtonGroup, QGridLayout)
from matplotlib.figure import Figure
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
from PyQt6.QtCore import Qt
from PyQt6.QtGui import QIcon, QColor
import pandas as pd

from datetime import datetime

from seaborn.external.docscrape import header

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import openpyxl
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment

from src.data.gendarmerie.structure import SUBDIVISIONS, SERVICE_RANGES


class FullListWindow(QMainWindow):
    """Fenêtre affichant la liste exhaustive des sanctionnés avec filtres."""

    def __init__(self, db_manager, parent=None):
        super().__init__(parent)
        self.filters = {}
        self.search_edit = None
        self.table = None
        self.result_label = None
        self.db_manager = db_manager
        self.setWindowTitle("Liste exhaustive des sanctionnés")
        self.setMinimumSize(1200, 800)

        self.setup_ui()
        self.load_filters()
        self.load_data()

    def setup_ui(self):
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)

        # Section Recherche
        search_group = QGroupBox("Recherche")
        search_layout = QVBoxLayout()

        # Radio buttons pour le type de recherche
        radio_layout = QHBoxLayout()
        self.search_type_group = QButtonGroup()

        self.matricule_radio = QRadioButton("Par Matricule")
        self.nom_radio = QRadioButton("Par Nom")
        self.matricule_radio.setChecked(True)

        self.search_type_group.addButton(self.matricule_radio)
        self.search_type_group.addButton(self.nom_radio)

        radio_layout.addWidget(self.matricule_radio)
        radio_layout.addWidget(self.nom_radio)
        radio_layout.addStretch()

        search_layout.addLayout(radio_layout)

        # Champ de recherche unique
        self.search_edit = QLineEdit()
        self.search_edit.setMinimumHeight(40)
        self.search_edit.setStyleSheet("""
            QLineEdit{
            border: 1px solid #fede77;
            border-radius: 5px;
            padding: 5px;
            }
            QLineEdit:focus{
            border: 1px solid #0078D4;
            }
        """)
        self.search_edit.setPlaceholderText("Entrez le matricule...")
        # Connecter la recherche dynamique
        self.search_edit.textChanged.connect(self.dynamic_search)
        search_layout.addWidget(self.search_edit)

        search_group.setLayout(search_layout)
        main_layout.addWidget(search_group)

        # Section Filtres
        filter_group = QGroupBox("Filtres")
        filter_layout = QGridLayout()

        # Configuration des filtres
        filters_config = [
            {"label": "Grade:", "widget": "combo", "name": "grade", "row": 0, "col": 0},
            {"label": "Subdivision:", "widget": "combo", "name": "subdiv", "row": 0, "col": 1},
            {"label": "Fautes commises:", "widget": "combo", "name": "faute", "row": 0, "col": 2},
            {"label": "Situation matrimoniale:", "widget": "combo", "name": "situation", "row": 0, "col": 3},
            {"label": "Tranche d'année de service:", "widget": "combo", "name": "service", "row": 1, "col": 0},
            {"label": "Année:", "widget": "combo", "name": "annee", "row": 1, "col": 1},
            {"label": "Statut:", "widget": "combo", "name": "statut", "row": 1, "col": 2},
            {"label": "Catégorie de fautes:", "widget": "combo", "name": "categorie", "row": 1, "col": 3}
        ]

        # Création des filtres
        self.filters = {}
        for filter_config in filters_config:
            # Layout vertical pour chaque filtre
            layout = QVBoxLayout()

            # Label
            label = QLabel(filter_config["label"])
            layout.addWidget(label)

            # ComboBox
            combo = QComboBox()
            combo.addItem(f"Tous(tes)")  # Option par défaut
            self.filters[filter_config["name"]] = combo
            layout.addWidget(combo)

            # Ajout au layout grille
            filter_layout.addLayout(layout, filter_config["row"], filter_config["col"])

        filter_group.setLayout(filter_layout)
        main_layout.addWidget(filter_group)

        # Boutons d'action
        button_layout = QHBoxLayout()
        self.apply_filter_btn = QPushButton("Appliquer les filtres")
        self.apply_filter_btn.setStyleSheet("""
                QPushButton {
                    background-color: #0A84FF;
                    color: white;
                    padding: 8px 15px;
                    border-radius: 5px;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #0066CC;
                }
            """)
        self.apply_filter_btn.clicked.connect(self.apply_filters)  # Connexion ajoutée

        self.reset_filter_btn = QPushButton("Réinitialiser")
        self.reset_filter_btn.setStyleSheet("""
                QPushButton {
                    background-color: #6C757D;
                    color: white;
                    padding: 8px 15px;
                    border-radius: 5px;
                    font-weight: bold;
                }
                QPushButton:hover {
                    background-color: #5A6268;
                }
            """)
        self.reset_filter_btn.clicked.connect(self.reset_filters)  # Connexion ajoutée

        button_layout.addWidget(self.apply_filter_btn)
        button_layout.addWidget(self.reset_filter_btn)
        button_layout.addStretch()

        main_layout.addLayout(button_layout)

        # Label pour le nombre de résultats
        self.result_label = QLabel()
        self.result_label.setStyleSheet("""
                QLabel {
                    font-size: 14px;
                    font-weight: bold;
                    color: #333;
                    padding: 5px;
                }
            """)
        main_layout.addWidget(self.result_label)

        # Tableau des données
        self.table = QTableWidget()
        self.table.setAlternatingRowColors(True)

        headers = ["ID", "Date d'enr", "Matricule", "Nom et Prénoms", "Grade", "Subdivision",
                   "Date des faits", "Faute commise", "Catégorie", "Statut",
                   "N° Dossier", "Années de service", "Situation matrimoniale"]
        self.table.setHorizontalHeaderLabels(headers)
        self.table.setSelectionBehavior(QTableWidget.SelectionBehavior.SelectRows)
        self.table.setEditTriggers(QTableWidget.EditTrigger.NoEditTriggers)
        main_layout.addWidget(self.table)

        # Boutons d'export
        export_layout = QHBoxLayout()

        self.btn_excel = QPushButton("Exporter Excel")
        self.btn_excel.setStyleSheet("""
            QPushButton {
                        background-color: #217346;
                        color: white;
                        padding: 8px 15px;
                        border-radius: 15px;
                        font-weight: bold;
                    }
                    QPushButton:hover {
                        background-color: #154734;
                    }
                    """)
        self.btn_excel.clicked.connect(self.export_excel)
        self.btn_excel.setIcon(QIcon("../resources/icons/excel_icon.svg"))

        self.btn_pdf = QPushButton("Exporter PDF")
        self.btn_pdf.setStyleSheet("""
                    QPushButton {
                                background-color: #6C63FF;
                                color: white;
                                padding: 8px 15px;
                                border-radius: 15px;
                                font-weight: bold;
                            }
                            QPushButton:hover {
                                background-color: #4B0082;
                            }
                            """)
        self.btn_pdf.setIcon(QIcon("../resources/icons/pdf.svg"))
        self.btn_pdf.clicked.connect(self.export_pdf)

        self.btn_pptx = QPushButton("Exporter PowerPoint")
        self.btn_pptx.setStyleSheet("""
                    QPushButton {
                                background-color: #D24726;
                                color: white;
                                padding: 8px 15px;
                                border-radius: 15px;
                                font-weight: bold;
                            }
                            QPushButton:hover {
                                background-color: #B7472A;
                            }
                            """)

        self.btn_pptx.setIcon(QIcon("../resources/icons/pptx.svg"))
        self.btn_pptx.clicked.connect(self.export_pptx)

        export_layout.addWidget(self.btn_excel)
        export_layout.addWidget(self.btn_pdf)
        export_layout.addWidget(self.btn_pptx)
        export_layout.addStretch()

        main_layout.addLayout(export_layout)

    def get_gendarme_info(self, matricule):
        """Récupère les informations d'un gendarme à partir de son matricule."""
        try:
            gendarme_query = """
            SELECT 
                nom_prenoms,
                grade,
                subdiv,
                annee_service,
                situation_matrimoniale
            FROM gendarmes 
            WHERE mle = ?
            """
            gendarme_params = [matricule]

            # Ajout des filtres spécifiques aux gendarmes
            if self.filters["grade"].currentText() != "Tous(tes)":
                gendarme_query += " AND grade = ?"
                gendarme_params.append(self.filters["grade"].currentText())

            if self.filters["subdiv"].currentText() != "Tous(tes)":
                gendarme_query += " AND subdiv = ?"
                gendarme_params.append(self.filters["subdiv"].currentText())

            if self.filters["situation"].currentText() != "Tous(tes)":
                gendarme_query += " AND situation_matrimoniale = ?"
                gendarme_params.append(self.filters["situation"].currentText())

            if self.filters["service"].currentText() != "Tous(tes)":
                service_text = self.filters["service"].currentText()
                # Extraire juste le nombre du texte (par exemple "15" de "15 ANS")
                try:
                    if "ANS" in service_text:
                        service_years = int(service_text.split()[0])
                        gendarme_query += " AND annee_service = ?"
                        gendarme_params.append(service_years)
                    elif "-" in service_text:  # Pour les tranches comme "0-5"
                        start, end = map(int, service_text.split("-"))
                        gendarme_query += " AND annee_service BETWEEN ? AND ?"
                        gendarme_params.extend([start, end])
                except ValueError as e:
                    print(f"Erreur de conversion de la tranche d'années: {e}")

            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute(gendarme_query, gendarme_params)
                result = cursor.fetchone()
                return result if result else ("", "", "", "", "")

        except Exception as e:
            print(f"Erreur dans get_gendarme_info: {str(e)}")
            return ("", "", "", "", "")

    def dynamic_search(self, text):
        """Effectue la recherche dynamique."""
        try:
            search_type = "matricule" if self.matricule_radio.isChecked() else "nom"

            # Si le champ est vide, afficher toutes les données
            if not text:
                self.load_data()
                return

            # Construction de la requête selon le type de recherche
            if search_type == "matricule":
                sanctions_query = """
                SELECT 
                    s.id,
                    s.date_enr,
                    s.matricule,
                    s.date_faits,
                    s.faute_commise,
                    s.categorie,
                    s.statut,
                    s.numero_dossier
                FROM sanctions s
                WHERE s.matricule LIKE ?
                ORDER BY s.id DESC
                """
                params = [f"%{text}%"]
            else:
                sanctions_query = """
                SELECT 
                    s.id,
                    s.date_enr,
                    s.matricule,
                    s.date_faits,
                    s.faute_commise,
                    s.categorie,
                    s.statut,
                    s.numero_dossier
                FROM sanctions s
                WHERE s.matricule IN (
                    SELECT mle FROM gendarmes 
                    WHERE nom_prenoms LIKE ?
                )
                ORDER BY s.id DESC
                """
                params = [f"%{text}%"]

            # Exécution de la requête
            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute(sanctions_query, params)
                sanctions_results = cursor.fetchall()

                # Mise à jour du tableau avec les résultats
                self.table.setRowCount(len(sanctions_results))

                filtered_count = 0
                for i, sanction in enumerate(sanctions_results):
                    # Récupérer les infos du gendarme
                    gendarme_info = self.get_gendarme_info(sanction[2])  # matricule est à l'index 2

                    # Création de la ligne complète
                    row_data = [
                        sanction[0],  # ID
                        self.format_date(sanction[1]),  # Date d'enr
                        sanction[2],  # Matricule
                        gendarme_info[0],  # Nom et Prénoms
                        gendarme_info[1],  # Grade
                        gendarme_info[2],  # Subdivision
                        self.format_date(sanction[3]),  # Date des faits
                        sanction[4],  # Faute commise
                        sanction[5],  # Catégorie
                        sanction[6],  # Statut
                        sanction[7],  # N° Dossier
                        gendarme_info[3],  # Années de service
                        gendarme_info[4]  # Situation Matrimoniale
                    ]

                    # Remplissage de la ligne
                    for j, value in enumerate(row_data):
                        item = QTableWidgetItem(str(value) if value is not None else "")

                        # Alignement à gauche pour nom/prénoms, faute et n° dossier
                        if j in [3, 7, 10]:
                            item.setTextAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
                        else:
                            item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)

                        # Coloration si le statut est "RADIE"
                        if sanction[6] == "RADIE":
                            item.setBackground(QColor(255, 200, 200))

                        self.table.setItem(filtered_count, j, item)

                    filtered_count += 1

                # Ajustement du nombre de lignes et mise à jour du label
                self.table.setRowCount(filtered_count)
                self.result_label.setText(f"Nombre de résultats : {filtered_count}")

        except Exception as e:
            print(f"Erreur dans la recherche dynamique : {str(e)}")  # Debug

    def update_table_with_results(self, sanctions_results):
        """Met à jour le tableau avec les résultats de la recherche."""
        self.table.setRowCount(len(sanctions_results))

        for i, sanction in enumerate(sanctions_results):
            # Récupérer les infos du gendarme
            gendarme_info = self.get_gendarme_info(sanction[1])  # matricule est à l'index 1

            # Même logique de remplissage que dans load_data
            row_data = [
                sanction[0],  # ID
                self.format_date(sanction[2]),  # Date d'enr
                sanction[1],  # Matricule
                gendarme_info[0],  # Nom et Prénoms
                gendarme_info[1],  # Grade
                gendarme_info[2],  # Subdivision
                self.format_date(sanction[4]),  # Date des faits
                sanction[3],  # Faute commise
                sanction[5],  # Catégorie
                sanction[6],  # Statut
                sanction[7],  # N° Dossier
                gendarme_info[3],  # Années de service
                gendarme_info[4]  # Situation Matrimoniale
            ]

            for j, value in enumerate(row_data):
                item = QTableWidgetItem(str(value) if value is not None else "")

                # Conserver les alignements spécifiques
                if j in [3, 7, 10]:
                    item.setTextAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
                else:
                    item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)

                # Conserver la coloration des radiés
                if sanction[6] == "RADIE":
                    item.setBackground(QColor(255, 200, 200))

                self.table.setItem(i, j, item)

    def load_filters(self):
        """Charge les valeurs des filtres."""
        try:
            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                print("\nDébug chargement des filtres:")
                # Grades
                cursor.execute("SELECT DISTINCT grade FROM gendarmes ORDER BY grade")
                grades = cursor.fetchall()
                self.filters["grade"].addItems([g[0] for g in grades if g[0]])
                print(f"Grades chargés: {[g[0] for g in grades if g[0]]}")

                # Subdivisions
                self.filters["subdiv"].addItems(SUBDIVISIONS)

                # Fautes commises
                cursor.execute("SELECT DISTINCT faute_commise FROM sanctions ORDER BY faute_commise")
                fautes = cursor.fetchall()
                self.filters["faute"].addItems([f[0] for f in fautes if f[0]])

                # Situation matrimoniale
                cursor.execute("SELECT DISTINCT situation_matrimoniale FROM gendarmes ORDER BY situation_matrimoniale")
                situations = cursor.fetchall()
                self.filters["situation"].addItems([s[0] for s in situations if s[0]])

                # Années
                cursor.execute("SELECT DISTINCT annee_punition FROM sanctions ORDER BY annee_punition DESC")
                annees = cursor.fetchall()
                self.filters["annee"].addItems([str(a[0]) for a in annees if a[0]])

                # Statuts
                cursor.execute("SELECT DISTINCT statut FROM sanctions ORDER BY statut")
                statuts = cursor.fetchall()
                self.filters["statut"].addItems([s[0] for s in statuts if s[0]])

                # Catégories
                cursor.execute("SELECT DISTINCT categorie FROM sanctions ORDER BY categorie")
                categories = cursor.fetchall()
                self.filters["categorie"].addItems([str(c[0]) for c in categories if c[0]])

                # Tranches d'années de service avec formatage
                formatted_ranges = [f"{range_text} ans" for range_text in SERVICE_RANGES]
                self.filters["service"].addItems(formatted_ranges)

        except Exception as e:
            QMessageBox.critical(self, "Erreur",
                                 f"Erreur lors du chargement des filtres: {str(e)}")

    def load_data(self):
        """Charge les données filtrées dans le tableau."""
        try:
            # 1. Requête pour les sanctions (données principales)
            sanctions_query = """
            SELECT 
                s.id,
                s.date_enr,
                s.matricule,
                s.date_faits,
                s.faute_commise,
                s.categorie,
                s.statut,
                s.numero_dossier
            FROM sanctions s
            WHERE 1=1
            """
            params = []

            # Application des filtres pour sanctions
            if self.filters["faute"].currentText() != "Tous(tes)":
                sanctions_query += " AND s.faute_commise = ?"
                params.append(self.filters["faute"].currentText())

            if self.filters["annee"].currentText() != "Tous(tes)":
                sanctions_query += " AND s.annee_punition = ?"
                params.append(int(self.filters["annee"].currentText()))

            if self.filters["statut"].currentText() != "Tous(tes)":
                sanctions_query += " AND s.statut = ?"
                params.append(self.filters["statut"].currentText())

            if self.filters["categorie"].currentText() != "Tous(tes)":
                sanctions_query += " AND s.categorie = ?"
                params.append(self.filters["categorie"].currentText())

            # Exécution de la requête principale
            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute(sanctions_query + " ORDER BY id DESC", params)
                sanctions_results = cursor.fetchall()

            # Configuration du tableau
            headers = ["ID", "Date d'enr", "Matricule", "Nom et Prénoms", "Grade", "Subdivision",
                       "Date des faits", "Faute commise", "Catégorie", "Statut",
                       "N° Dossier", "Années de service", "Situation Matrimoniale"]

            self.table.setRowCount(len(sanctions_results))
            self.table.setColumnCount(len(headers))
            self.table.setHorizontalHeaderLabels(headers)

            # Remplissage des données
            filtered_count = 0
            for i, sanction in enumerate(sanctions_results):
                # Récupérer les infos du gendarme avec filtres
                gendarme_info = self.get_gendarme_info(sanction[2])  # Utilisation de la méthode de classe

                # Si les infos du gendarme sont vides à cause des filtres, continuer
                if not any(gendarme_info):
                    continue

                # Création de la ligne complète
                row_data = [
                    sanction[0],  # ID
                    self.format_date(sanction[1]),  # Date d'enr
                    sanction[2],  # Matricule
                    gendarme_info[0],  # Nom et Prénoms
                    gendarme_info[1],  # Grade
                    gendarme_info[2],  # Subdivision
                    self.format_date(sanction[3]),  # Date des faits
                    sanction[4],  # Faute commise
                    sanction[5],  # Catégorie
                    sanction[6],  # Statut
                    sanction[7],  # N° Dossier
                    gendarme_info[3],  # Années de service
                    gendarme_info[4]  # Situation Matrimoniale
                ]

                # Remplissage de la ligne avec alignements et coloration
                for j, value in enumerate(row_data):
                    item = QTableWidgetItem(str(value) if value is not None else "")

                    # Alignement à gauche pour nom/prénoms, faute et n° dossier
                    if j in [3, 7, 10]:
                        item.setTextAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
                    else:
                        item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)

                    # Coloration si le statut est "RADIE"
                    if sanction[6] == "RADIE":
                        item.setBackground(QColor(255, 200, 200))

                    self.table.setItem(filtered_count, j, item)

                filtered_count += 1

            # Ajuster le nombre de lignes au nombre réel après filtrage
            self.table.setRowCount(filtered_count)

            # Ajustement des colonnes
            self.table.horizontalHeader().setSectionResizeMode(
                QHeaderView.ResizeMode.ResizeToContents
            )

            # Mise à jour du label de résultats
            self.result_label.setText(f"Nombre de résultats : {filtered_count}")

        except Exception as e:
            print(f"Erreur dans load_data: {str(e)}")  # Pour le debug
            QMessageBox.critical(self, "Erreur",
                                 f"Erreur lors du chargement des données: {str(e)}")

    def format_date(self, date_str):
        """Formate une date en JJ/MM/AAAA."""
        if date_str:
            try:
                date = datetime.strptime(str(date_str), "%Y-%m-%d")
                return date.strftime("%d/%m/%Y")
            except:
                return str(date_str)
        return ""

    def apply_filters(self):
        """Applique les filtres sélectionnés."""
        print("\nDébug application des filtres:")
        try:
            # Debug: afficher les valeurs des filtres
            print("Valeurs des filtres:")
            for key, combo in self.filters.items():
                print(f"{key}: {combo.currentText()}")

            self.load_data()  # Cette méthode contient déjà la logique de filtrage

        except Exception as e:
            print(f"Erreur dans apply_filters: {str(e)}")
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'application des filtres: {str(e)}"
            )

    def load_data_with_filters(self, query, params):
        """Charge les données avec les filtres appliqués."""
        try:
            with self.db_manager.get_connection() as conn:
                cursor = conn.cursor()
                cursor.execute(query + " ORDER BY id DESC", params)
                sanctions_results = cursor.fetchall()

            # Configuration du tableau
            headers = ["ID", "Date d'enr", "Matricule", "Nom et Prénoms", "Grade", "Subdivision",
                       "Date des faits", "Faute commise", "Catégorie", "Statut",
                       "N° Dossier", "Années de service", "Situation Matrimoniale"]

            self.table.setRowCount(len(sanctions_results))
            self.table.setColumnCount(len(headers))
            self.table.setHorizontalHeaderLabels(headers)

            # Remplissage des données
            filtered_count = 0
            for i, sanction in enumerate(sanctions_results):
                gendarme_info = self.get_gendarme_info(sanction[2])

                # Si les filtres sur gendarme excluent cet enregistrement, passer
                if not any(gendarme_info):
                    continue

                row_data = [
                    sanction[0],  # ID
                    self.format_date(sanction[1]),  # Date d'enr
                    sanction[2],  # Matricule
                    gendarme_info[0],  # Nom et Prénoms
                    gendarme_info[1],  # Grade
                    gendarme_info[2],  # Subdivision
                    self.format_date(sanction[3]),  # Date des faits
                    sanction[4],  # Faute commise
                    sanction[5],  # Catégorie
                    sanction[6],  # Statut
                    sanction[7],  # N° Dossier
                    gendarme_info[3],  # Années de service
                    gendarme_info[4]  # Situation Matrimoniale
                ]

                # Remplissage avec mise en forme
                for j, value in enumerate(row_data):
                    item = QTableWidgetItem(str(value) if value is not None else "")

                    if j in [3, 7, 10]:
                        item.setTextAlignment(Qt.AlignmentFlag.AlignLeft | Qt.AlignmentFlag.AlignVCenter)
                    else:
                        item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)

                    if sanction[6] == "RADIE":
                        item.setBackground(QColor(255, 200, 200))

                    self.table.setItem(filtered_count, j, item)

                filtered_count += 1

            # Ajuster la taille finale
            self.table.setRowCount(filtered_count)

            # Ajuster les colonnes
            self.table.horizontalHeader().setSectionResizeMode(
                QHeaderView.ResizeMode.ResizeToContents
            )

            # Mise à jour du compteur
            self.result_label.setText(f"Nombre de résultats : {filtered_count}")

        except Exception as e:
            print(f"Erreur dans load_data_with_filters: {str(e)}")
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors du chargement des données filtrées: {str(e)}"
            )

    def reset_filters(self):
        """Réinitialise tous les filtres."""
        print("\nRéinitialisation des filtres")
        for combo in self.filters.values():
            combo.setCurrentIndex(0)  # Remet à "Tous(tes)"
        self.search_edit.clear()  # Efface aussi la recherche
        self.load_data()

    def export_excel(self):
        """Exporte les données vers Excel avec mise en forme."""
        try:
            # Demander à l'utilisateur où sauvegarder le fichier
            default_name = f"liste_sanctions_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer le fichier Excel",
                default_name,
                "Excel Files (*.xlsx)"
            )

            if not file_path:  # L'utilisateur a annulé.
                return

            # Création d'un writer Excel
            with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
                # Récupération des données du tableau
                data = []
                headers = []

                # Récupération des en-têtes
                for j in range(self.table.columnCount()):
                    headers.append(self.table.horizontalHeaderItem(j).text())

                # Récupération des données
                for i in range(self.table.rowCount()):
                    row = []
                    for j in range(self.table.columnCount()):
                        item = self.table.item(i, j)
                        row.append(item.text() if item else "")
                    data.append(row)

                # Création du DataFrame
                df = pd.DataFrame(data, columns=headers)

                # Export vers Excel avec mise en forme
                df.to_excel(writer, sheet_name='Liste des sanctions', index=False)
                worksheet = writer.sheets['Liste des sanctions']

                # Ajustement des colonnes
                for idx, col in enumerate(df.columns):
                    max_length = max(
                        df[col].astype(str).apply(len).max(),
                        len(str(col))
                    ) + 2
                    worksheet.column_dimensions[openpyxl.utils.get_column_letter(idx + 1)].width = max_length

                # Style pour l'en-tête
                for cell in worksheet[1]:
                    cell.font = openpyxl.styles.Font(bold=True)
                    cell.fill = openpyxl.styles.PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")

                # Bordures et alignement
                thin_border = openpyxl.styles.borders.Border(
                    left=openpyxl.styles.borders.Side(style='thin'),
                    right=openpyxl.styles.borders.Side(style='thin'),
                    top=openpyxl.styles.borders.Side(style='thin'),
                    bottom=openpyxl.styles.borders.Side(style='thin')
                )

                for row in worksheet.iter_rows(min_row=1, max_row=len(data) + 1):
                    for cell in row:
                        cell.border = thin_border
                        cell.alignment = openpyxl.styles.Alignment(horizontal='center', vertical='center')

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export Excel:\n{str(e)}"
            )

    def export_pdf(self):
        """Exporte les données vers PDF avec mise en forme."""
        try:
            # Demander à l'utilisateur où sauvegarder le fichier
            default_name = f"liste_sanctions_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer le fichier PDF",
                default_name,
                "PDF Files (*.pdf)"
            )

            if not file_path:  # L'utilisateur a annulé
                return

            # Création du document PDF
            doc = SimpleDocTemplate(
                file_path,
                pagesize=landscape(letter),
                rightMargin=30,
                leftMargin=30,
                topMargin=30,
                bottomMargin=30
            )

            # Liste des éléments à ajouter au PDF
            elements = []
            styles = getSampleStyleSheet()

            # Titre du document
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=30,
                alignment=1  # Centre
            )
            title = Paragraph("Liste des Sanctions", title_style)
            elements.append(title)

            # Informations sur les filtres appliqués
            filter_info = []
            if self.filters["grade"].currentText() != "Tous(tes)":
                filter_info.append(f"Grade: {self.filters['grade'].currentText()}")
            if self.filters["subdiv"].currentText() != "Tous(tes)":
                filter_info.append(f"Subdivision: {self.filters['subdiv'].currentText()}")
            if self.filters["faute"].currentText() != "Tous(tes)":
                filter_info.append(f"Faute commise: {self.filters['faute'].currentText()}")
            if self.filters["situation"].currentText() != "Tous(tes)":
                filter_info.append(f"Situation matrimoniale: {self.filters['situation'].currentText()}")
            if self.filters["annee"].currentText() != "Tous(tes)":
                filter_info.append(f"Année de punition: {self.filters['annee'].currentText()}")
            if self.filters["statut"].currentText() != "Tous(tes)":
                filter_info.append(f"Statut: {self.filters['statut'].currentText()}")
            if self.filters["categorie"].currentText() != "Tous(tes)":
                filter_info.append(f"Catégorie: {self.filters['categorie'].currentText()}")
            if self.filters["service"].currentText() != "Tous(tes)":
                filter_info.append(f"Années de service: {self.filters['service'].currentText()}")

            if filter_info:
                filter_text = "Filtres appliqués: " + ", ".join(filter_info)
                elements.append(Paragraph(filter_text, styles['Normal']))
                elements.append(Spacer(1, 20))

            # Création des données du tableau
            table_data = []
            headers = []

            # En-têtes
            for j in range(self.table.columnCount()):
                headers.append(self.table.horizontalHeaderItem(j).text())
            table_data.append(headers)

            # Données
            for i in range(self.table.rowCount()):
                row = []
                for j in range(self.table.columnCount()):
                    item = self.table.item(i, j)
                    row.append(item.text() if item else "")
                table_data.append(row)

            # Style du tableau
            table_style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('ROWHEIGHT', (0, 0), (-1, -1), 20)
            ])

            pdf_table = Table(table_data, repeatRows=1)
            pdf_table.setStyle(table_style)
            elements.append(pdf_table)

            # Informations de pagination
            def add_page_number(canvas, mydoc):
                page_num = canvas.getPageNumber()
                text = f"Page {page_num}"
                canvas.saveState()
                canvas.setFont('Helvetica', 8)
                canvas.drawRightString(mydoc.pagesize[0] - 30, 30, text)
                canvas.restoreState()

            # Construction du document
            doc.build(elements, onFirstPage=add_page_number, onLaterPages=add_page_number)

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export PDF:\n{str(e)}"
            )

    def export_pptx(self):
        """Exporte les données vers PowerPoint avec mise en forme."""
        try:
            # Demander à l'utilisateur où sauvegarder le fichier
            default_name = f"liste_sanctions_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            file_path, _ = QFileDialog.getSaveFileName(
                self,
                "Enregistrer la présentation PowerPoint",
                default_name,
                "PowerPoint Files (*.pptx)"
            )

            if not file_path:  # L'utilisateur a annulé
                return

            # Création de la présentation
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Slide de titre
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "Liste des Sanctions"
            subtitle.text = f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}"

            # Slide d'information sur les filtres
            bullet_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = "Filtres appliqués"

            tf = body_shape.text_frame

            # Informations sur les filtres appliqués
            filter_info = []
            if self.filters["grade"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Grade: {self.filters['grade'].currentText()}"

            if self.filters["subdiv"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Subdivision: {self.filters['subdiv'].currentText()}"

            if self.filters["faute"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Faute commise: {self.filters['faute'].currentText()}"

            if self.filters["situation"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Situation matrimoniale: {self.filters['situation'].currentText()}"

            if self.filters["annee"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Année de punition: {self.filters['annee'].currentText()}"

            if self.filters["statut"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Statut: {self.filters['statut'].currentText()}"

            if self.filters["categorie"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Categorie de faute: {self.filters['categorie'].currentText()}"

            if self.filters["service"].currentText() != "Tous(tes)":
                p = tf.add_paragraph()
                p.text = f"Année de service: {self.filters['service'].currentText()}"

            # Slides de données
            rows_per_slide = 10
            total_rows = self.table.rowCount()
            num_slides = (total_rows + rows_per_slide - 1) // rows_per_slide

            headers = [self.table.horizontalHeaderItem(j).text()
                       for j in range(self.table.columnCount())]

            for slide_num in range(num_slides):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes

                title = shapes.title
                title.text = f"Liste des Sanctions ({slide_num + 1}/{num_slides})"

                # Indices de début et fin pour ce slide
                start_idx = slide_num * rows_per_slide
                end_idx = min((slide_num + 1) * rows_per_slide, total_rows)

                # Création du tableau
                rows = end_idx - start_idx + 1
                cols = self.table.columnCount()

                left = Inches(0.5)
                top = Inches(1.5)
                width = prs.slide_width - Inches(1)
                height = Inches(0.3 * rows)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # Style de l'en-tête
                for i, header in enumerate(headers):
                    cell = table.cell(0, i)
                    cell.text = header
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(128, 128, 128)

                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)

                # Données
                for i in range(start_idx, end_idx):
                    for j in range(cols):
                        cell = table.cell(i - start_idx + 1, j)
                        item = self.table.item(i, j)
                        cell.text = item.text() if item else ""

                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(9)
                        paragraph.alignment = PP_ALIGN.CENTER

            # Sauvegarde de la présentation
            prs.save(file_path)

            QMessageBox.information(
                self,
                "Succès",
                f"Les données ont été exportées vers:\n{file_path}"
            )

        except Exception as e:
            QMessageBox.critical(
                self,
                "Erreur",
                f"Erreur lors de l'export PowerPoint:\n{str(e)}"
            )
